from pptx import Presentation

prs = Presentation(r'C:\Users\Asus\Desktop\PeerReview_Presentation.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Width: {prs.slide_width.inches:.2f}"  Height: {prs.slide_height.inches:.2f}"')
for i, sl in enumerate(prs.slides):
    texts = []
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()[:80].replace('\n', ' ')
            if t:
                texts.append(t)
    print(f'Slide {i+1}: {" | ".join(texts[:3])}')
