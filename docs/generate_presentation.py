#!/usr/bin/env python3
"""
Academic Peer Review Evaluation System — Polished Minimal Presentation
Design: Clean white cards, single blue accent, generous whitespace.
        No shadows, no badge pills, no mixed colors — just clear hierarchy.
Run:    python docs/generate_presentation.py
Out:    C:/Users/Asus/Desktop/PeerReview_Presentation_v2.pptx
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Palette (strict: dark + blue + muted, nothing else) ─────────────────────
SLATE     = RGBColor(248, 250, 252)   # slide background
WHITE     = RGBColor(255, 255, 255)   # card background
DARK      = RGBColor(15,  23,  42)    # headings
BODY      = RGBColor(51,  65,  85)    # body text
MUTED     = RGBColor(148, 163, 184)   # secondary / labels
BLUE      = RGBColor(37,  99,  235)   # primary accent
BLUE_SOFT = RGBColor(239, 246, 255)   # tinted bg for callouts
BORDER    = RGBColor(226, 232, 240)   # card border
CODE_BG   = RGBColor(248, 250, 252)   # code background
CODE_FG   = RGBColor(30,  41,  59)    # code text
NAV_BG    = RGBColor(15,  23,  42)    # top nav bar
SUCCESS   = RGBColor(5,   150, 105)   # used ONLY for checkmarks
AMBER     = RGBColor(217, 119,  6)    # used ONLY for warnings

LOGO = (
    r"C:\Users\Asus\.cursor\projects\d-Peer-Review-Evaluation-System"
    r"\assets\c__Users_Asus_AppData_Roaming_Cursor_User_workspaceStorage"
    r"_56a4694182cc1675762ac69cecfe39bd_images_LogoGIFblue_transp_600w"
    r"-7aca3195-955a-4bd3-b33c-1240e323c5a9.png"
)
OUT = r"C:\Users\Asus\Desktop\PeerReview_Presentation_v2.pptx"

# ─── Presentation ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ─── Layout constants ─────────────────────────────────────────────────────────
NAV_H = Inches(0.38)     # nav bar
TY    = Inches(0.46)     # title Y
TH    = Inches(0.52)     # title box height
DIV_Y = TY + TH          # divider line   ≈ 0.98"
CY    = DIV_Y + Inches(0.14)  # content Y ≈ 1.12"
CX    = Inches(0.52)
CW    = Inches(12.29)
CH    = H - CY - Inches(0.2)  # content height ≈ 6.18"
GAP   = Inches(0.22)     # gap between panels

SECTIONS = [
    "Motivation", "System Design", "Technical Setup",
    "Data & Modeling", "Testing", "Live System", "Future Work",
]

# slide counter helper
_slide_idx = [0]
def _next():
    _slide_idx[0] += 1
    return _slide_idx[0]

# ─── Animation XML ────────────────────────────────────────────────────────────
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

_TIMING_XML = '''\
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
    <p:childTnLst><p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>{G}</p:childTnLst></p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
    </p:seq></p:childTnLst>
  </p:cTn></p:par></p:tnLst>
  <p:bldLst>{B}</p:bldLst>
</p:timing>'''

_GRP_XML   = '<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{I}" fill="hold"><p:stCondLst><p:cond delay="{T}"/></p:stCondLst><p:childTnLst>{S}</p:childTnLst></p:cTn></p:par>'
_FADE_XML  = '<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{A}" fill="hold"><p:stCondLst><p:cond delay="{D}"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{B}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="{G}" nodeType="{N}"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{C}" dur="{U}"/><p:tgtEl><p:spTgt spid="{S}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
_BLD_XML   = '<p:bldP xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spid="{S}" grpId="{G}" build="p"/>'


def animate(slide, groups):
    sl = slide._element
    ex = sl.find(f'{{{_P}}}timing')
    if ex is not None:
        sl.remove(ex)
    cid = [3]
    def ni(): cid[0] += 1; return cid[0]
    all_g, all_b = [], []
    for gi, g in enumerate(groups, 1):
        shapes  = g['shapes']
        stagger = g.get('stagger', 120)
        dur     = g.get('dur', 380)
        auto    = g.get('auto', False)
        s_xml = []
        for i, sh in enumerate(shapes):
            nt = ('afterEffect' if auto else 'clickEffect') if i == 0 else 'withEffect'
            s_xml.append(_FADE_XML.format(A=ni(), B=ni(), C=ni(), D=i*stagger,
                                          G=gi, N=nt, U=dur, S=sh.shape_id))
            all_b.append(_BLD_XML.format(S=sh.shape_id, G=gi))
        all_g.append(_GRP_XML.format(I=ni(), T='0' if auto else 'indefinite',
                                     S=''.join(s_xml)))
    sl.append(etree.fromstring(_TIMING_XML.format(
        G=''.join(all_g), B=''.join(all_b)).encode()))


# ─── Drawing primitives ───────────────────────────────────────────────────────
def I(x): return int(x)


def rect(sl, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    sh = sl.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = I(lw)
    else:
        sh.line.width = 0
    return sh


def tb(sl, x, y, w, h, text='', sz=Pt(13), col=BODY,
       bold=False, ital=False, align=PP_ALIGN.LEFT,
       font='Calibri', wrap=True):
    box = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = sz; run.font.color.rgb = col
    run.font.bold = bold; run.font.italic = ital; run.font.name = font
    return box, tf


def para(tf, text, sz=Pt(13), col=BODY, bold=False, ital=False,
         align=PP_ALIGN.LEFT, font='Calibri', sp=Pt(5)):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = sp
    run = p.add_run()
    run.text = text; run.font.size = sz; run.font.color.rgb = col
    run.font.bold = bold; run.font.italic = ital; run.font.name = font
    return p


def logo(sl, x=None, y=Inches(0.0), w=Inches(1.32), h=Inches(0.38)):
    if os.path.exists(LOGO):
        lx = x if x is not None else W - w - Inches(0.1)
        sl.shapes.add_picture(LOGO, I(lx), I(y), I(w), I(h))


# ─── Panel (the only card style used) ────────────────────────────────────────
def panel(sl, x, y, w, h, heading=None, lines=None,
          accent=BLUE, h_sz=Pt(15), b_sz=Pt(13), bg=WHITE):
    """
    Clean panel: white bg, 1px border, optional heading with small
    blue underline accent, body lines below.  No shadows, no filled headers.
    """
    rect(sl, x, y, w, h, fill=bg, line=BORDER, lw=Pt(0.75))
    body_y = y + Inches(0.16)
    if heading:
        tb(sl, x + Inches(0.22), y + Inches(0.14),
           w - Inches(0.32), Inches(0.34),
           text=heading, sz=h_sz, col=DARK, bold=True)
        rect(sl, x + Inches(0.22), y + Inches(0.52), Inches(0.3), Pt(2.5), fill=accent)
        body_y = y + Inches(0.62)
    if lines:
        bh = y + h - body_y - Inches(0.16)
        if isinstance(lines, str):
            tb(sl, x + Inches(0.22), body_y, w - Inches(0.32), bh,
               text=lines, sz=b_sz, col=BODY)
        else:
            _, tf = tb(sl, x + Inches(0.22), body_y, w - Inches(0.32), bh,
                       text=lines[0], sz=b_sz, col=BODY)
            for ln in lines[1:]:
                para(tf, ln, sz=b_sz, sp=Pt(6))
    return sl.shapes[-1] if lines else sl.shapes[-2]


def code_block(sl, lines, x, y, w, h, sz=Pt(9.5)):
    """Clean GitHub-style code block."""
    rect(sl, x, y, w, h, fill=CODE_BG, line=BORDER, lw=Pt(0.75))
    box = sl.shapes.add_textbox(I(x + Inches(0.18)), I(y + Inches(0.16)),
                                I(w - Inches(0.36)), I(h - Inches(0.28)))
    tf = box.text_frame; tf.word_wrap = False
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0.5)
        run = p.add_run()
        run.text = ln; run.font.size = sz
        run.font.color.rgb = CODE_FG; run.font.name = 'Consolas'
    return box


def ph(sl, label, x, y, w, h):
    """Diagram placeholder."""
    rect(sl, x, y, w, h, fill=BLUE_SOFT, line=BORDER, lw=Pt(1.0))
    rect(sl, x, y + h/2 - Inches(0.48), w, Pt(1), fill=BORDER)
    rect(sl, x, y + h/2 - Inches(0.14), w, Pt(1), fill=BORDER)
    rect(sl, x, y + h/2 + Inches(0.2), w, Pt(1), fill=BORDER)
    tb(sl, x, y + h/2 - Inches(0.75), w, Inches(0.44),
       text=f'Insert from Report PDF: {label}',
       sz=Pt(13), col=BLUE, bold=True, ital=True, align=PP_ALIGN.CENTER)
    tb(sl, x, y + h/2 + Inches(0.35), w, Inches(0.34),
       text='( copy & paste diagram here )',
       sz=Pt(11), col=MUTED, ital=True, align=PP_ALIGN.CENTER)


# ─── Slide factory ────────────────────────────────────────────────────────────
def new_slide(section, title):
    n = _next()
    sl = prs.slides.add_slide(BLANK)
    # SLATE background
    rect(sl, 0, 0, W, H, fill=SLATE)
    # White content area
    rect(sl, 0, NAV_H, W, H - NAV_H, fill=WHITE)
    # Nav bar
    rect(sl, 0, 0, W, NAV_H, fill=NAV_BG)
    tb(sl, Inches(0.2), Inches(0.06), Inches(0.85), NAV_H - Inches(0.1),
       text=f'{n} / 18', sz=Pt(9), col=MUTED, align=PP_ALIGN.CENTER)
    tb(sl, Inches(1.1), Inches(0.06), W - Inches(2.6), NAV_H - Inches(0.1),
       text=section.upper(), sz=Pt(10), col=WHITE, bold=True,
       align=PP_ALIGN.CENTER, wrap=False)
    logo(sl, y=Inches(0.0), w=Inches(1.28), h=Inches(0.38))
    # Title
    tb(sl, CX, TY, W - Inches(1.0), TH,
       text=title, sz=Pt(24), col=DARK, bold=True, wrap=False)
    rect(sl, CX, DIV_Y, CW, Pt(1.5), fill=BLUE)
    return sl


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=WHITE)
rect(sl, 0, 0, W, Inches(0.05), fill=BLUE)   # top accent stripe

# Left hero content
tb(sl, Inches(0.6), Inches(0.55), Inches(7.6), Inches(0.3),
   text='ACADEMIC PLATFORM', sz=Pt(10.5), col=BLUE, bold=True)

tb(sl, Inches(0.6), Inches(0.9), Inches(7.6), Inches(1.65),
   text='Academic Peer Review\nEvaluation System',
   sz=Pt(42), col=DARK, bold=True)

rect(sl, Inches(0.6), Inches(2.68), Inches(0.32), Pt(2), fill=BLUE)

tb(sl, Inches(0.6), Inches(2.8), Inches(6.8), Inches(0.8),
   text='A modular platform that takes students from assignment submission\n'
        'through anonymous peer review to final graded results — all in one place.',
   sz=Pt(13.5), col=MUTED)

# Info row
tb(sl, Inches(0.6), Inches(3.72), Inches(6.8), Inches(0.36),
   text='Tennessee State University  ·  Software Systems Design  ·  April 2026',
   sz=Pt(12), col=MUTED, ital=True)

# TSU Logo (left bottom area)
if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, I(Inches(0.6)), I(Inches(4.22)),
                          I(Inches(2.6)), I(Inches(1.0)))

# Right: 2×2 clean cards
RX1 = Inches(7.95); RY1 = Inches(0.5)
CRW = Inches(2.58); CRH = Inches(1.5); CRG = Inches(0.18)
feat = [
    ('Courses & Enrollments',
     'Instructors create courses;\nstudents enroll and access\nassignments.'),
    ('File Submissions',
     'Students upload work before\nthe deadline with optional\nnotes.'),
    ('Anonymous Review',
     'Reviewers score rubric\ncriteria without seeing\nauthor identity.'),
    ('Grading & Feedback',
     'Final grades computed from\nreview averages, released\nwhen ready.'),
]
title_shapes = []
for i, (ttl, body) in enumerate(feat):
    col_i, row_i = i % 2, i // 2
    fx = RX1 + col_i * (CRW + CRG)
    fy = RY1 + row_i * (CRH + CRG)
    sh = panel(sl, fx, fy, CRW, CRH, ttl, body, h_sz=Pt(12.5), b_sz=Pt(11))
    title_shapes.append(sh)

# Workflow strip
rect(sl, 0, Inches(5.95), W, Inches(1.55), fill=SLATE, line=None)
rect(sl, 0, Inches(5.95), W, Pt(1), fill=BORDER)
tb(sl, 0, Inches(6.07), W, Inches(0.3),
   text='WORKFLOW', sz=Pt(10), col=MUTED, bold=True, align=PP_ALIGN.CENTER)

wf = ['Create assignment', 'Student submits', 'Allocate reviewers',
      'Peer reviews', 'Grade & release']
total_wf = len(wf) * Inches(2.1) + (len(wf)-1) * Inches(0.35)
wfx0 = (W - total_wf) / 2
for i, step in enumerate(wf):
    sx = wfx0 + i * (Inches(2.1) + Inches(0.35))
    tb(sl, sx, Inches(6.48), Inches(2.1), Inches(0.38),
       text=step, sz=Pt(12), col=BODY, bold=True, align=PP_ALIGN.CENTER)
    if i < len(wf) - 1:
        tb(sl, sx + Inches(2.1) + Inches(0.02), Inches(6.48),
           Inches(0.33), Inches(0.38),
           text='›', sz=Pt(16), col=MUTED, align=PP_ALIGN.CENTER)

animate(sl, [{'shapes': title_shapes, 'auto': True, 'stagger': 140, 'dur': 380}])


# =============================================================================
# SLIDE 2 — Outline
# =============================================================================
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=WHITE)
rect(sl, 0, 0, W, Inches(0.05), fill=BLUE)
rect(sl, 0, Inches(0.05), W, Inches(0.82), fill=NAV_BG)
logo(sl, y=Inches(0.09), w=Inches(1.32), h=Inches(0.62))
tb(sl, Inches(0.5), Inches(0.18), Inches(10.5), Inches(0.58),
   text='Presentation Outline', sz=Pt(28), col=WHITE, bold=True)

outline = [
    ('01', 'Motivation',       'Problem statement, requirements'),
    ('02', 'System Design',    'Tech stack, app map, SDLC, layered architecture'),
    ('03', 'Technical Setup',  'Docker, PostgreSQL, packages, Django ORM & Admin'),
    ('04', 'Data & Modeling',  'ER model, class diagram, sequence diagrams'),
    ('05', 'Testing',          'Verification, validation & testing strategy'),
    ('06', 'Live System',      'Application UI walkthrough'),
    ('07', 'Future Work',      'Planned improvements & next steps'),
]
OW = Inches(5.78); OH = Inches(0.72); OG = Inches(0.14)
OX0 = Inches(0.52); OX1 = Inches(6.98); OY0 = Inches(1.08)

all_cards = []
for i, (num, sec, sub) in enumerate(outline):
    xi = OX0 if i < 4 else OX1
    yi = OY0 + (i if i < 4 else i - 4) * (OH + OG)
    rect(sl, xi, yi, OW, OH, fill=WHITE, line=BORDER, lw=Pt(0.75))
    # Number
    tb(sl, xi + Inches(0.16), yi + Inches(0.1), Inches(0.5), OH - Inches(0.16),
       text=num, sz=Pt(15), col=BLUE, bold=True, font='Calibri')
    # Divider
    rect(sl, xi + Inches(0.72), yi + Inches(0.1), Pt(1), OH - Inches(0.18), fill=BORDER)
    # Section
    tb(sl, xi + Inches(0.88), yi + Inches(0.1), OW - Inches(1.0), Inches(0.3),
       text=sec, sz=Pt(14), col=DARK, bold=True)
    tb(sl, xi + Inches(0.88), yi + Inches(0.42), OW - Inches(1.0), Inches(0.26),
       text=sub, sz=Pt(11), col=MUTED, ital=True)
    all_cards.append(sl.shapes[-1])

animate(sl, [{'shapes': all_cards, 'auto': True, 'stagger': 80, 'dur': 320}])


# =============================================================================
# SLIDE 3 — Motivation  (Section: Motivation)
# =============================================================================
sl = new_slide('Motivation', 'The Problem with Traditional Peer Review')

LW3 = (CW - GAP) / 2  # ~6.04"
RX3 = CX + LW3 + GAP

# Left panel — challenges
panel(sl, CX, CY, LW3, CH, heading='The Challenge')
probs = [
    '  Manual review is inconsistent at scale',
    '  Reviewer bias when identities are visible',
    '  No structured rubric-to-score traceability',
    '  Error-prone grade computation from multiple reviewers',
    '  No real-time review completion visibility',
    '  No audit trail per rubric criterion',
]
_, tp = tb(sl, CX + Inches(0.22), CY + Inches(0.7),
           LW3 - Inches(0.32), CH - Inches(0.78),
           text=probs[0], sz=Pt(14), col=BODY)
for p in probs[1:]:
    para(tp, p, sz=Pt(14), col=BODY, sp=Pt(11))

# Right panel — solution
panel(sl, RX3, CY, LW3, CH, heading='Our Solution', accent=SUCCESS)
sols = [
    '  Full lifecycle: enrollment to result release',
    '  UUID tokens — identity never revealed',
    '  Rubric + per-criterion scoring',
    '  Automated normalisation → mean → GPA',
    '  Role dashboards with live stats',
    '  CALCULATED → RELEASED grade lifecycle',
]
_, ts = tb(sl, RX3 + Inches(0.22), CY + Inches(0.7),
           LW3 - Inches(0.32), CH - Inches(0.78),
           text=sols[0], sz=Pt(14), col=BODY)
for s in sols[1:]:
    para(ts, s, sz=Pt(14), col=BODY, sp=Pt(11))


# =============================================================================
# SLIDE 4 — Tech Stack  (Section: System Design)
# =============================================================================
sl = new_slide('System Design', 'Tech Stack')

stack = [
    ('Backend',  'Django 5.x',          'Class-based views, service layer, ORM'),
    ('Database', 'PostgreSQL 15',        'psycopg2-binary · ACID · Django ORM'),
    ('Frontend', 'Bootstrap 5 + Icons', 'CDN-only · no React/Vue · Responsive'),
    ('Auth',     'Custom User Model',   'AbstractUser + role (Admin/Instructor/Student)'),
    ('Config',   'python-dotenv',       '.env pattern — secrets outside settings.py'),
    ('Media',    'Django FileField',    'Local media/ · path: assignment_<id>/<reg>/'),
]
CW4 = (CW - 2*GAP) / 3  # ~3.95"
CH4 = (CH - GAP) / 2    # ~2.98"

card_shapes = []
for i, (layer, tech, desc) in enumerate(stack):
    col4, row4 = i % 3, i // 3
    x4 = CX + col4 * (CW4 + GAP)
    y4 = CY + row4  * (CH4 + GAP)
    rect(sl, x4, y4, CW4, CH4, fill=WHITE, line=BORDER, lw=Pt(0.75))
    # Layer label (small caps, muted)
    tb(sl, x4 + Inches(0.22), y4 + Inches(0.18),
       CW4 - Inches(0.3), Inches(0.28),
       text=layer.upper(), sz=Pt(9.5), col=MUTED, bold=True)
    # Technology name (big, blue)
    tb(sl, x4 + Inches(0.22), y4 + Inches(0.5),
       CW4 - Inches(0.3), Inches(0.58),
       text=tech, sz=Pt(17), col=BLUE, bold=True)
    # Description
    tb(sl, x4 + Inches(0.22), y4 + Inches(1.12),
       CW4 - Inches(0.3), Inches(1.55),
       text=desc, sz=Pt(13), col=MUTED)
    card_shapes.append(sl.shapes[-1])

animate(sl, [{'shapes': card_shapes, 'auto': True, 'stagger': 70, 'dur': 300}])


# =============================================================================
# SLIDE 5 — Project Structure  (Section: System Design)
# =============================================================================
sl = new_slide('System Design', 'Project Structure — 7 Django Apps')

apps = [
    ('accounts',    'Custom User, role gate mixin, role dashboards'),
    ('courses',     'Course CRUD + Enrollment management'),
    ('assignments', 'Assignment CRUD, course-scoped, deadline tracking'),
    ('submissions', 'File upload, one submission per student per assignment'),
    ('rubrics',     'Rubric + RubricCriterion, one rubric per assignment'),
    ('reviews',     'Reviewer allocation service, anonymous UUID review forms'),
    ('grading',     'Grade calculation + release, FinalGrade + ReviewerAccuracy'),
]
# config banner
rect(sl, (W - Inches(3.2)) / 2, CY, Inches(3.2), Inches(0.44), fill=NAV_BG)
tb(sl, (W - Inches(3.2)) / 2, CY, Inches(3.2), Inches(0.44),
   text='config  (settings · urls · wsgi)',
   sz=Pt(12), col=WHITE, bold=True, align=PP_ALIGN.CENTER, font='Consolas')

AH = Inches(0.72); AG = Inches(0.11)
AW = Inches(5.85)
AY0 = CY + Inches(0.58)

app_shapes = []
for i, (name, desc) in enumerate(apps):
    col_a = i % 2; row_a = i // 2
    if i < 4:
        ax, ay = CX, AY0 + i * (AH + AG)
    else:
        ax, ay = CX + AW + GAP, AY0 + (i - 4) * (AH + AG)
    aw = AW
    rect(sl, ax, ay, aw, AH, fill=WHITE, line=BORDER, lw=Pt(0.75))
    # app name (monospace)
    tb(sl, ax + Inches(0.18), ay + Inches(0.1), Inches(1.45), AH - Inches(0.16),
       text=name, sz=Pt(13), col=BLUE, bold=True, font='Consolas', wrap=False)
    # Separator
    rect(sl, ax + Inches(1.7), ay + Inches(0.1), Pt(1), AH - Inches(0.18), fill=BORDER)
    # Description
    tb(sl, ax + Inches(1.85), ay + Inches(0.14), aw - Inches(2.0), AH - Inches(0.22),
       text=desc, sz=Pt(12), col=BODY)
    app_shapes.append(sl.shapes[-1])

animate(sl, [{'shapes': app_shapes, 'auto': True, 'stagger': 80, 'dur': 300}])


# =============================================================================
# SLIDE 6 — PostgreSQL & Docker  (Section: Technical Setup)
# =============================================================================
sl = new_slide('Technical Setup', 'Database Setup — PostgreSQL & Docker')

# Subtitle
tb(sl, CX, CY, CW, Inches(0.3),
   text='PostgreSQL 15 in Docker — isolated, persistent, cleanly configured.',
   sz=Pt(13), col=MUTED, ital=True)

CODE_W = Inches(7.5); CODE_H = Inches(3.15)
code_block(sl, [
    '# Step 1 — Named volume (data persists across restarts)',
    'docker volume create peer_review_volume',
    '',
    '# Step 2 — Launch container',
    'docker run --name peer-review-pg \\',
    '  --env  POSTGRES_DB=peer_review_db       \\',
    '  --env  POSTGRES_USER=peer_user          \\',
    '  --env  POSTGRES_PASSWORD=securepass     \\',
    '  --publish 5433:5432                     \\   # 5432 was taken locally',
    '  --volume peer_review_volume:/var/lib/postgresql/data \\',
    '  --detach postgres:15',
], CX, CY + Inches(0.36), CODE_W, CODE_H)

code_block(sl, [
    '# Step 3 — Create DB user',
    'docker exec -it peer-review-pg psql -U postgres',
    '',
    'CREATE USER peer_user',
    '    WITH PASSWORD \'securepass\';',
    'GRANT ALL PRIVILEGES ON DATABASE',
    '    peer_review_db TO peer_user;',
    '',
    '# Step 4 — Migrate',
    'python manage.py migrate',
], Inches(8.15), CY + Inches(0.36), CW - CODE_W - GAP, CODE_H)

# Two callout panels below
PW6 = (CW - GAP) / 2

panel(sl, CX, CY + Inches(0.36) + CODE_H + Inches(0.18), PW6, Inches(2.42),
      heading='Port Conflict — Challenge & Fix',
      lines=[
          'Port 5432 was already occupied by a local PostgreSQL installation.',
          '',
          'Fix: --publish 5433:5432 maps Docker port 5432 → host 5433.',
          'Django .env:  DJANGO_DB_PORT=5433',
      ], accent=AMBER, b_sz=Pt(12.5))

panel(sl, CX + PW6 + GAP, CY + Inches(0.36) + CODE_H + Inches(0.18), PW6, Inches(2.42),
      heading='Named Volume — Data Persistence',
      lines=[
          'peer_review_volume mounts to /var/lib/postgresql/data.',
          '',
          'Database survives container stop & restart.',
          'Data lifecycle fully decoupled from container.',
      ], accent=SUCCESS, b_sz=Pt(12.5))


# =============================================================================
# SLIDE 7 — Packages & Config  (Section: Technical Setup)
# =============================================================================
sl = new_slide('Technical Setup', 'Python Packages & Django Configuration')

code_block(sl, [
    '# requirements.txt',
    'django>=5.0',
    'psycopg2-binary',
    'python-dotenv',
], CX, CY, Inches(3.72), Inches(1.55), sz=Pt(11))

code_block(sl, [
    '# .env  (git-ignored)',
    'DJANGO_SECRET_KEY=<key>',
    'DJANGO_DEBUG=True',
    '',
    'DJANGO_DB_NAME=peer_review_db',
    'DJANGO_DB_USER=peer_user',
    'DJANGO_DB_PASSWORD=securepass',
    'DJANGO_DB_HOST=localhost',
    'DJANGO_DB_PORT=5433',
], Inches(4.06), CY, Inches(4.38), Inches(2.65))

code_block(sl, [
    '# config/settings.py',
    'from dotenv import load_dotenv',
    'load_dotenv(BASE_DIR / \'.env\')',
    '',
    'AUTH_USER_MODEL = \'accounts.User\'',
    '',
    'DATABASES = {',
    '  \'default\': {',
    '    \'ENGINE\':   \'django.db.backends.postgresql\',',
    '    \'NAME\':     os.getenv(\'DJANGO_DB_NAME\'),',
    '    \'USER\':     os.getenv(\'DJANGO_DB_USER\'),',
    '    \'PASSWORD\': os.getenv(\'DJANGO_DB_PASSWORD\'),',
    '    \'HOST\':     os.getenv(\'DJANGO_DB_HOST\', \'localhost\'),',
    '    \'PORT\':     os.getenv(\'DJANGO_DB_PORT\', \'5432\'),',
    '  }',
    '}',
    '',
    'MEDIA_URL  = \'media/\'',
    'MEDIA_ROOT = BASE_DIR / \'media\'',
], CX, CY + Inches(1.72), Inches(8.28), Inches(4.38))

# Package notes (right column)
pkgs = [
    ('django>=5.0',       'CBVs, ORM, Admin, Auth, Migrations'),
    ('psycopg2-binary',   'PostgreSQL adapter; binary — no libpq required'),
    ('python-dotenv',     'Loads .env into os.environ at startup'),
    ('Bootstrap 5 (CDN)', 'Responsive UI & icons — zero build step'),
]
PX7 = Inches(8.65); py7 = CY
for pkg, note in pkgs:
    PH7 = Inches(1.46); PW7 = CW - (PX7 - CX)
    rect(sl, PX7, py7, PW7, PH7, fill=WHITE, line=BORDER, lw=Pt(0.75))
    tb(sl, PX7 + Inches(0.18), py7 + Inches(0.14),
       PW7 - Inches(0.26), Inches(0.32),
       text=pkg, sz=Pt(13), col=DARK, bold=True, font='Consolas')
    rect(sl, PX7 + Inches(0.18), py7 + Inches(0.5), Inches(0.28), Pt(2), fill=BLUE)
    tb(sl, PX7 + Inches(0.18), py7 + Inches(0.6),
       PW7 - Inches(0.26), PH7 - Inches(0.7),
       text=note, sz=Pt(12.5), col=MUTED)
    py7 += PH7 + Inches(0.14)


# =============================================================================
# SLIDE 8 — Django ORM  (Section: Technical Setup)
# =============================================================================
sl = new_slide('Technical Setup', 'Django ORM in Action')

LW8 = Inches(6.08); RX8 = CX + LW8 + GAP; RW8 = CW - LW8 - GAP

code_block(sl, [
    '# Reviewer allocation — workload-balanced',
    'eligible = (',
    '    User.objects',
    '    .filter(',
    '        enrollments__course=assignment.course,',
    '        role=User.Roles.STUDENT,',
    '    )',
    '    .exclude(submissions__assignment=assignment)',
    '    .annotate(workload=Count(\'review_assignments\'))',
    '    .order_by(\'workload\')    # fewest reviews first',
    ')',
], CX, CY, LW8, Inches(2.75))

code_block(sl, [
    '# Grade calculation',
    'reviews    = Review.objects.filter(',
    '    review_assignment__submission=submission,',
    '    review_assignment__review_status=\'COMPLETED\',',
    ')',
    'normalised = [r.total_score / rubric_max * 100',
    '              for r in reviews]',
    'final      = clamp(mean(normalised), 0, 100)',
], CX, CY + Inches(2.9), LW8, Inches(2.15) + (CH - Inches(2.9) - Inches(2.15)))

# Right: key ORM patterns
panel(sl, RX8, CY, RW8, CH, heading='Key ORM Patterns')
kp = [
    'select_related()',
    '  → SQL JOIN, prevents N+1 queries',
    '',
    'annotate()',
    '  → workload count per reviewer',
    '',
    'exclude()',
    '  → no self-review possible',
    '',
    'OneToOneField',
    '  → one rubric per assignment',
    '',
    'UniqueConstraint',
    '  → one submission per student',
]
_, ktf = tb(sl, RX8 + Inches(0.22), CY + Inches(0.7),
            RW8 - Inches(0.32), CH - Inches(0.8),
            text=kp[0], sz=Pt(12.5), col=BODY, font='Consolas')
for k in kp[1:]:
    c = MUTED if k.startswith('  ') or k == '' else BODY
    para(ktf, k, sz=Pt(12.5), col=c, font='Consolas', sp=Pt(3))


# =============================================================================
# SLIDE 9 — Django Admin  (Section: Technical Setup)
# =============================================================================
sl = new_slide('Technical Setup', 'Django Admin — Full Model Visibility')

code_block(sl, [
    '# submissions/admin.py',
    '@admin.register(Submission)',
    'class SubmissionAdmin(admin.ModelAdmin):',
    '    list_display  = [\'student\',\'assignment\',\'course\',',
    '                     \'status\',\'has_file\',\'submitted_at\']',
    '    list_filter   = [\'status\',\'assignment__course\']',
    '    search_fields = [\'student__username\',\'student__reg_no\']',
    '',
    '    @admin.display(boolean=True)',
    '    def has_file(self, obj): return bool(obj.file)',
    '',
    '# grading/admin.py',
    '@admin.register(FinalGrade)',
    'class FinalGradeAdmin(admin.ModelAdmin):',
    '    list_display = [\'student\',\'assignment\',\'course\',',
    '                    \'numeric_score_100\',\'letter_grade\',',
    '                    \'gpa\',\'grade_status\',\'released_at\']',
    '    list_filter  = [\'grade_status\',',
    '                    \'submission__assignment__course\']',
], CX, CY, Inches(7.3), CH)

# Admin table (right)
TX9 = Inches(8.0); TW9 = CW - (TX9 - CX)
tbl_rows = [
    ('UserAdmin',             'Password hashing, reg_no fieldsets'),
    ('CourseAdmin',           'Enrollment count, instructor name'),
    ('AssignmentAdmin',       'Submission count, rubric bool'),
    ('SubmissionAdmin',       'File link, has_file bool'),
    ('RubricAdmin',           'RubricCriterionInline (tabular)'),
    ('ReviewAssignmentAdmin', 'Author + assignment + course'),
    ('ReviewAdmin',           'ReviewCriterionScoreInline'),
    ('FinalGradeAdmin',       'Score, grade, GPA, status'),
    ('ReviewerAccuracyAdmin', 'Accuracy + deviation display'),
]
HDR9 = Inches(0.4); ROW9 = Inches(0.59)
rect(sl, TX9, CY, TW9, HDR9, fill=NAV_BG)
tb(sl, TX9 + Inches(0.14), CY + Inches(0.07), Inches(2.12), HDR9 - Inches(0.1),
   text='Admin Class', sz=Pt(11), col=WHITE, bold=True)
tb(sl, TX9 + Inches(2.3), CY + Inches(0.07), TW9 - Inches(2.4), HDR9 - Inches(0.1),
   text='Key Features', sz=Pt(11), col=WHITE, bold=True)
for i, (cls, feat) in enumerate(tbl_rows):
    ry = CY + HDR9 + i * ROW9
    rect(sl, TX9, ry, TW9, ROW9,
         fill=BLUE_SOFT if i % 2 == 0 else WHITE, line=BORDER, lw=Pt(0.5))
    tb(sl, TX9 + Inches(0.12), ry + Inches(0.1), Inches(2.12), ROW9 - Inches(0.14),
       text=cls, sz=Pt(9.5), col=BLUE, bold=True, font='Consolas')
    tb(sl, TX9 + Inches(2.3), ry + Inches(0.1), TW9 - Inches(2.4), ROW9 - Inches(0.14),
       text=feat, sz=Pt(11.5), col=BODY)


# =============================================================================
# SLIDE 10 — Workflow  (Section: Technical Setup)
# =============================================================================
sl = new_slide('Technical Setup', 'End-to-End System Workflow')

steps = [
    ('1', 'Admin\nSetup',     'Admin creates users\nvia Django Admin panel'),
    ('2', 'Course\n& Enroll', 'Instructor creates course\n& enrolls students'),
    ('3', 'Assignment',       'Instructor creates\nassignment + due date'),
    ('4', 'Submission',       'Student uploads one\nfile per assignment'),
    ('5', 'Rubric',           'Instructor defines rubric\n+ weighted criteria'),
    ('6', 'Allocation',       'System assigns 3\nanonymous reviewers'),
    ('7', 'Review',           'Reviewers score each\ncriterion + comments'),
    ('8', 'Grading',          'Normalise → average\n→ letter grade + GPA'),
    ('9', 'Release',          'Instructor releases;\nstudent sees feedback'),
]
SW = Inches(1.26); SH = Inches(5.55); SG = Inches(0.065)
TOTAL_S = len(steps)*SW + (len(steps)-1)*SG
SX0 = (W - TOTAL_S) / 2
SY0 = CY + Inches(0.06)

step_shapes = []
for i, (num, lbl, desc) in enumerate(steps):
    sx = SX0 + i*(SW+SG)
    sh = rect(sl, sx, SY0, SW, SH, fill=WHITE, line=BORDER, lw=Pt(0.75))
    step_shapes.append(sh)
    # Number circle
    rect(sl, sx + SW/2 - Inches(0.2), SY0 + Inches(0.14),
         Inches(0.4), Inches(0.4), fill=BLUE)
    tb(sl, sx + SW/2 - Inches(0.2), SY0 + Inches(0.14),
       Inches(0.4), Inches(0.4),
       text=num, sz=Pt(14), col=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Label
    tb(sl, sx + Inches(0.06), SY0 + Inches(0.62),
       SW - Inches(0.1), Inches(0.9),
       text=lbl, sz=Pt(10.5), col=DARK, bold=True)
    # Description
    tb(sl, sx + Inches(0.06), SY0 + Inches(1.6),
       SW - Inches(0.1), Inches(3.7),
       text=desc, sz=Pt(9.5), col=MUTED)
    # Thin connector
    if i < len(steps)-1:
        rect(sl, sx+SW, SY0 + SH/2 - Inches(0.03), SG, Inches(0.06), fill=BORDER)

animate(sl, [{'shapes': [sh], 'auto': False, 'dur': 280} for sh in step_shapes])


# =============================================================================
# SLIDE 11 — Grade Calculation  (Section: Data & Modeling)
# =============================================================================
sl = new_slide('Data & Modeling', 'Grade Calculation & Reviewer Accuracy')

CODE_W11 = Inches(7.68)
code_block(sl, [
    '# grading/services.py',
    'def calculate_grades_for_assignment(assignment, force=False):',
    '    rubric_max = rubric.total_max_marks()',
    '    for sub in assignment.submissions.all():',
    '        reviews = completed_reviews(sub)',
    '        if len(reviews) < REVIEWS_PER_SUBMISSION:',
    '            continue                 # not enough reviews yet',
    '',
    '        # Normalise to 0-100',
    '        normalised = [r.total_score / rubric_max * 100 for r in reviews]',
    '        final = clamp(mean(normalised), 0, 100)',
    '        letter, gpa = get_grade(final)  # grade_scale.py lookup',
    '',
    '        FinalGrade.objects.update_or_create(',
    '            submission=sub,',
    '            defaults=dict(numeric_score_100=final,',
    '                          letter_grade=letter, gpa=gpa))',
    '',
    '        for rev in reviews:',
    '            dev = abs(rev_normalised - final)',
    '            ReviewerAccuracy.objects.update_or_create(',
    '                review=rev,',
    '                defaults=dict(deviation=dev,',
    '                              accuracy=clamp(100-dev, 0, 100)))',
], CX, CY, CODE_W11, CH)

# Right: grade scale + accuracy
RX11 = CX + CODE_W11 + GAP; RW11 = CW - CODE_W11 - GAP
# Grade scale table
GH11 = Inches(0.5)
rect(sl, RX11, CY, RW11, Inches(0.38), fill=NAV_BG)
tb(sl, RX11 + Inches(0.14), CY + Inches(0.06), RW11 - Inches(0.2), Inches(0.28),
   text='Grade Scale', sz=Pt(12), col=WHITE, bold=True)
grade_rows = [('≥ 90','A','4.00'), ('≥ 80','B','3.00'), ('≥ 70','C','2.00'),
              ('≥ 60','D','1.00'), ('< 60','F','0.00')]
for i, (sc, lt, gpa_v) in enumerate(grade_rows):
    gy = CY + Inches(0.38) + i*GH11
    rect(sl, RX11, gy, RW11, GH11,
         fill=BLUE_SOFT if i%2==0 else WHITE, line=BORDER, lw=Pt(0.5))
    tb(sl, RX11+Inches(0.14), gy+Inches(0.1), Inches(1.1), GH11-Inches(0.12),
       text=sc, sz=Pt(13), col=BODY, font='Consolas')
    lc = SUCCESS if lt=='A' else (BLUE if lt in 'BC' else (AMBER if lt=='D' else
          RGBColor(180,40,40)))
    tb(sl, RX11+Inches(1.4), gy+Inches(0.08), Inches(0.5), GH11-Inches(0.1),
       text=lt, sz=Pt(15), col=lc, bold=True)
    tb(sl, RX11+Inches(2.1), gy+Inches(0.1), Inches(1.3), GH11-Inches(0.12),
       text=gpa_v+' GPA', sz=Pt(12.5), col=MUTED)

ACY = CY + Inches(0.38) + 5*GH11 + Inches(0.16)
panel(sl, RX11, ACY, RW11, CH - (ACY - CY), heading='Reviewer Accuracy')
acc_lines = [
    'deviation = |reviewer_norm − final_avg|',
    'accuracy  = clamp(100 − deviation, 0, 100)',
    '',
    '100  →  perfectly on average',
    '  0  →  maximally deviant',
    '',
    'Stored in ReviewerAccuracy.',
    'Instructor views per assignment.',
]
_, atf = tb(sl, RX11+Inches(0.22), ACY+Inches(0.7),
            RW11-Inches(0.32), CH-(ACY-CY)-Inches(0.8),
            text=acc_lines[0], sz=Pt(11), col=BODY, font='Consolas')
for a in acc_lines[1:]:
    c = BODY if '=' in a else MUTED
    f = 'Consolas' if '=' in a or '→' in a else 'Calibri'
    para(atf, a, sz=Pt(11), col=c, font=f, sp=Pt(4))


# =============================================================================
# SLIDE 12 — ER Diagram  (Section: Data & Modeling)
# =============================================================================
sl = new_slide('Data & Modeling', 'Entity-Relationship Diagram')

ph(sl, 'ER Diagram', CX, CY, Inches(8.08), CH)

RX12 = CX + Inches(8.08) + GAP; RW12 = CW - Inches(8.08) - GAP
panel(sl, RX12, CY, RW12, CH, heading='Key Relationships')
rels = [
    'User 1─N Course',
    'Course 1─N Enrollment',
    'Course 1─N Assignment',
    'Assignment 1─N Submission',
    'Assignment 1─1 Rubric',
    'Rubric 1─N RubricCriterion',
    'Submission 1─N ReviewAssignment',
    'ReviewAssignment 1─1 Review',
    'Review 1─N ReviewCriterionScore',
    'Submission 1─1 FinalGrade',
    'Review 1─1 ReviewerAccuracy',
]
_, rtf = tb(sl, RX12+Inches(0.22), CY+Inches(0.7),
            RW12-Inches(0.32), CH-Inches(0.8),
            text=rels[0], sz=Pt(11), col=BODY, font='Consolas')
for rel in rels[1:]:
    para(rtf, rel, sz=Pt(11), col=BODY, font='Consolas', sp=Pt(6))


# =============================================================================
# SLIDE 13 — Class + Sequence Diagrams  (Section: Data & Modeling)
# =============================================================================
sl = new_slide('Data & Modeling', 'Structural & Interaction Models')

HW13 = (CW - GAP) / 2
ph(sl, 'Class Diagram', CX, CY, HW13, CH)
ph(sl, 'Sequence Diagram — Submit Assignment Flow',
   CX + HW13 + GAP, CY, HW13, CH)

# Labels
rect(sl, CX, CY, HW13, Inches(0.32), fill=NAV_BG)
tb(sl, CX + Inches(0.14), CY + Inches(0.05), HW13 - Inches(0.2), Inches(0.24),
   text='Class Diagram', sz=Pt(11), col=WHITE, bold=True)
rect(sl, CX + HW13 + GAP, CY, HW13, Inches(0.32), fill=BLUE)
tb(sl, CX + HW13 + GAP + Inches(0.14), CY + Inches(0.05), HW13 - Inches(0.2), Inches(0.24),
   text='Sequence Diagram — Submit Assignment Flow', sz=Pt(11), col=WHITE, bold=True)


# =============================================================================
# SLIDE 14 — V&V / Testing  (Section: Testing)
# =============================================================================
sl = new_slide('Testing', 'Verification, Validation & Testing')

LW14 = (CW - GAP) / 2; RX14 = CX + LW14 + GAP

panel(sl, CX, CY, LW14, CH, heading='Verification')
ver = [
    ('Code Reviews',        'Thin views; services.py pattern enforced throughout'),
    ('Django Guardrails',   'model.clean() validates every save before DB write'),
    ('DB Constraints',      'UniqueConstraint + CheckConstraint at PostgreSQL level'),
    ('Migration Tracking',  'All schema changes tracked as files in git'),
    ('Role Gate Coverage',  'RoleRequiredMixin applied to every non-public view'),
]
vy = CY + Inches(0.7)
for hd, bd in ver:
    tb(sl, CX + Inches(0.22), vy, LW14 - Inches(0.32), Inches(0.26),
       text=hd, sz=Pt(13.5), col=DARK, bold=True)
    tb(sl, CX + Inches(0.22), vy + Inches(0.3), LW14 - Inches(0.32), Inches(0.42),
       text=bd, sz=Pt(12.5), col=MUTED)
    vy += Inches(0.88)

panel(sl, RX14, CY, LW14, CH, heading='Validation', accent=SUCCESS)
val = [
    ('Role-Based Testing',  'Each role tested for correct access + redirect'),
    ('Ownership Checks',    'Instructors cannot access other instructors\' data'),
    ('Anonymous Review',    'Reviewer/author identity never exposed in context'),
    ('Allocation Rules',    'Self-review blocked; 3 reviews per submission enforced'),
    ('Grade Lifecycle',     'CALCULATED → RELEASED; unrelease path tested'),
    ('Edge Cases',          'Duplicate submission blocked; late status correct'),
]
vay = CY + Inches(0.7)
for hd, bd in val:
    tb(sl, RX14 + Inches(0.22), vay, LW14 - Inches(0.32), Inches(0.26),
       text=hd, sz=Pt(13.5), col=DARK, bold=True)
    tb(sl, RX14 + Inches(0.22), vay + Inches(0.3), LW14 - Inches(0.32), Inches(0.42),
       text=bd, sz=Pt(12.5), col=MUTED)
    vay += Inches(0.88)


# =============================================================================
# SLIDE 15 — Live System  (Section: Live System)
# =============================================================================
sl = new_slide('Live System', 'Live System — Application UI')

tb(sl, CX, CY, CW, Inches(0.28),
   text='Django 5 + Bootstrap 5 — deployed locally, fully functional.',
   sz=Pt(13), col=MUTED, ital=True)

screens = [
    'Admin Dashboard',       'Instructor Dashboard',
    'Student Dashboard',     'Submission Upload',
    'Anonymous Review Form', 'Grade Result View',
]
SW15 = (CW - 2*GAP) / 3; SH15 = (CH - GAP - Inches(0.32)) / 2
for i, sc in enumerate(screens):
    col15, row15 = i%3, i//3
    sx15 = CX + col15*(SW15+GAP)
    sy15 = CY + Inches(0.34) + row15*(SH15+GAP)
    ph(sl, sc, sx15, sy15, SW15, SH15)
    rect(sl, sx15, sy15, SW15, Inches(0.3), fill=NAV_BG if row15==0 else BLUE)
    tb(sl, sx15+Inches(0.12), sy15+Inches(0.05), SW15-Inches(0.2), Inches(0.22),
       text=sc, sz=Pt(10.5), col=WHITE, bold=True)


# =============================================================================
# SLIDE 16 — Future Work  (Section: Future Work)
# =============================================================================
sl = new_slide('Future Work', 'Future Improvements')

future = [
    ('Email Notifications',    'Alert students when reviews are due or grades released'),
    ('Rubric Templates',       'Clone and reuse rubrics across assignments'),
    ('Anti-Collusion Checks',  'Detect suspiciously aligned scores within peer groups'),
    ('Review Analytics',       'Score distribution charts per assignment'),
    ('REST API / Mobile',      'Expose endpoints for a future mobile or SPA client'),
    ('Production Hardening',   'HTTPS, S3 media, split settings, structured logging'),
]
CW16 = (CW - 2*GAP) / 3; CH16 = (CH - GAP) / 2
f_shapes = []
for i, (ttl, desc) in enumerate(future):
    col16, row16 = i%3, i//3
    fx = CX + col16*(CW16+GAP)
    fy = CY + row16*(CH16+GAP)
    rect(sl, fx, fy, CW16, CH16, fill=WHITE, line=BORDER, lw=Pt(0.75))
    tb(sl, fx+Inches(0.22), fy+Inches(0.18), CW16-Inches(0.32), Inches(0.3),
       text=ttl, sz=Pt(15), col=DARK, bold=True)
    rect(sl, fx+Inches(0.22), fy+Inches(0.52), Inches(0.3), Pt(2.5), fill=BLUE)
    tb(sl, fx+Inches(0.22), fy+Inches(0.64), CW16-Inches(0.32), CH16-Inches(0.78),
       text=desc, sz=Pt(13), col=MUTED)
    f_shapes.append(sl.shapes[-1])

animate(sl, [{'shapes': f_shapes, 'auto': True, 'stagger': 90, 'dur': 300}])


# =============================================================================
# SLIDE 17 — SDLC & Agile  (Section: System Design)
# =============================================================================
# insert at position 6 (after app map) — but python-pptx appends, so we just add it
# We note it's numbered dynamically anyway
sl = new_slide('System Design', 'SDLC & Agile Practices')

LW17 = (CW - GAP) / 2; RX17 = CX + LW17 + GAP

panel(sl, CX, CY, LW17, CH, heading='Incremental Development (Sommerville)')
phases = [
    ('Increment 1', 'accounts + courses + assignments'),
    ('Increment 2', 'submissions + rubrics'),
    ('Increment 3', 'reviews'),
    ('Increment 4', 'grading'),
]
py17 = CY + Inches(0.7)
for inc, apps_str in phases:
    rect(sl, CX + Inches(0.22), py17, LW17 - Inches(0.32), Inches(1.02),
         fill=BLUE_SOFT, line=BORDER, lw=Pt(0.5))
    tb(sl, CX + Inches(0.36), py17 + Inches(0.1),
       Inches(1.3), Inches(0.28), text=inc, sz=Pt(12), col=BLUE, bold=True)
    tb(sl, CX + Inches(0.36), py17 + Inches(0.42),
       LW17 - Inches(0.6), Inches(0.44),
       text=apps_str, sz=Pt(12.5), col=BODY, font='Consolas')
    py17 += Inches(1.16)

panel(sl, RX17, CY, LW17, CH, heading='Agile Practices Integrated')
agile = [
    ('Iterative Delivery',  'Working software after each increment'),
    ('XP — Simple Design',  'Thin views; service-layer business logic'),
    ('XP — Refactoring',    'services.py extracted iteratively from views'),
    ('Small Releases',      'Each app independently committed + demo-ed'),
    ('Continuous Testing',  'Manual regression after each increment'),
    ('Customer Collab.',    'User stories drove each increment scope'),
]
ay17 = CY + Inches(0.7)
for hd, bd in agile:
    tb(sl, RX17 + Inches(0.22), ay17, LW17 - Inches(0.32), Inches(0.28),
       text=hd, sz=Pt(13.5), col=DARK, bold=True)
    tb(sl, RX17 + Inches(0.22), ay17 + Inches(0.32), LW17 - Inches(0.32), Inches(0.38),
       text=bd, sz=Pt(12.5), col=MUTED)
    ay17 += Inches(0.85)


# =============================================================================
# SLIDE 18 — Any Questions?
# =============================================================================
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=WHITE)
rect(sl, 0, 0, W, Inches(0.05), fill=BLUE)
rect(sl, 0, H - Inches(0.05), W, Inches(0.05), fill=BLUE)
# Right accent panel
rect(sl, Inches(8.8), 0, Inches(4.53), H, fill=SLATE)
rect(sl, Inches(8.8), 0, Pt(1.5), H, fill=BORDER)

# TSU Logo
if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, I(Inches(5.45)), I(Inches(0.42)),
                          I(Inches(2.4)), I(Inches(0.9)))

# Main message
tb(sl, Inches(0.6), Inches(1.6), Inches(7.8), Inches(0.28),
   text='OPEN DISCUSSION', sz=Pt(10.5), col=BLUE, bold=True)
tb(sl, Inches(0.6), Inches(1.98), Inches(7.8), Inches(1.2),
   text='Thank you!', sz=Pt(50), col=DARK, bold=True, wrap=False)
rect(sl, Inches(0.6), Inches(3.25), Inches(5.2), Pt(1.5), fill=BORDER)
tb(sl, Inches(0.6), Inches(3.4), Inches(7.8), Inches(0.9),
   text='Any Questions?', sz=Pt(36), col=BLUE, bold=True, wrap=False)
tb(sl, Inches(0.6), Inches(4.45), Inches(7.8), Inches(0.5),
   text='Academic Peer Review Evaluation System  ·  Tennessee State University  ·  April 2026',
   sz=Pt(12), col=MUTED, ital=True)

# Right panel summary
items_q = [
    ('7 Apps',         'Fully modular Django structure'),
    ('12 Models',      'DB-enforced constraints throughout'),
    ('3 Roles',        'Admin, Instructor, Student'),
    ('Anonymous',      'UUID tokens, zero identity leakage'),
    ('Auto Grading',   'Normalise → average → GPA mapping'),
]
ry_q = Inches(0.65)
for key, val_q in items_q:
    tb(sl, Inches(9.1), ry_q, Inches(1.1), Inches(0.3),
       text=key, sz=Pt(11), col=BLUE, bold=True)
    tb(sl, Inches(10.28), ry_q, Inches(2.85), Inches(0.3),
       text=val_q, sz=Pt(11.5), col=MUTED)
    rect(sl, Inches(9.1), ry_q + Inches(0.34), Inches(4.0), Pt(0.75), fill=BORDER)
    ry_q += Inches(1.22)

# Animate
q_shapes = [sl.shapes[i] for i in range(2, min(7, len(sl.shapes)))]
animate(sl, [{'shapes': q_shapes, 'auto': True, 'stagger': 160, 'dur': 420}])


# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'\nSaved ({len(prs.slides)} slides, polished minimal) -> {OUT}\n')
