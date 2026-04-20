from pptx import Presentation
import os

path = r'C:\Users\Asus\Desktop\PeerReview_Presentation_v2.pptx'
prs = Presentation(path)
size_kb = os.path.getsize(path) // 1024
print(f'File size: {size_kb} KB')
print(f'Total slides: {len(prs.slides)}')
print(f'Dimensions: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
for i, sl in enumerate(prs.slides):
    texts = []
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()[:65].replace('\n', ' ')
            if t:
                texts.append(t)
    first_two = '  |  '.join(texts[:2])
    print(f'Slide {i+1:>2}: {first_two}')
